#!/usr/bin/env python3
"""OCR -> editable PPTX builder.

The script runs OCR one image at a time, caches each page as JSON, and then
rebuilds editable text layers over the clean slide backgrounds.
"""
from __future__ import annotations

import argparse
import importlib.util
import json
import os
import platform
import statistics
import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


SCRIPT_DIR = Path(__file__).resolve().parent

SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000


def px_to_emu_x(px: float, image_w: float) -> int:
    return int(round(px / image_w * SLIDE_W_EMU))


def px_to_emu_y(px: float, image_h: float) -> int:
    return int(round(px / image_h * SLIDE_H_EMU))


def default_font_name() -> str:
    if platform.system() == "Darwin":
        return "PingFang SC"
    if platform.system() == "Windows":
        return "Microsoft YaHei"
    return "Noto Sans CJK SC"


def set_ea_font(run, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    latin = rpr.get_or_add_latin()
    latin.set("typeface", font_name)
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", font_name)


def color_hex(color: tuple[int, int, int]) -> str:
    return "".join(f"{max(0, min(255, int(v))):02X}" for v in color)


def color_distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return sum((a[i] - b[i]) ** 2 for i in range(3)) ** 0.5


def median_color(pixels: list[tuple[int, int, int]]) -> tuple[int, int, int]:
    return tuple(int(round(statistics.median(channel))) for channel in zip(*pixels))


def mix_color(
    color: tuple[int, int, int],
    target: tuple[int, int, int],
    amount: float,
) -> tuple[int, int, int]:
    return tuple(
        int(round(color[i] * (1.0 - amount) + target[i] * amount))
        for i in range(3)
    )


def brighten_text_color(color: tuple[int, int, int], bg_brightness: float) -> tuple[int, int, int]:
    """Recover glyph-core brightness lost to antialias/background blending."""
    brightness = sum(color)
    saturation = max(color) - min(color)
    bg_is_dark = bg_brightness < 420
    if not bg_is_dark:
        return color

    if saturation < 38:
        if brightness < 520:
            return mix_color(color, (238, 246, 255), 0.72)
        if brightness < 650:
            return mix_color(color, (238, 246, 255), 0.42)
        return color

    if brightness < 390:
        return mix_color(color, (70, 185, 255), 0.48)
    if brightness < 510:
        return mix_color(color, (92, 205, 255), 0.32)
    return color


def core_ink_color(pixels: list[tuple[int, int, int]]) -> tuple[int, int, int]:
    """Use the bright glyph core, not antialiased edge pixels."""
    if not pixels:
        return (235, 246, 255)

    by_brightness = sorted(pixels, key=sum)
    if len(by_brightness) >= 10:
        start = int(len(by_brightness) * 0.58)
        end = max(start + 1, int(len(by_brightness) * 0.96))
        bright_core = by_brightness[start:end]
    else:
        bright_core = by_brightness

    # If cyan/blue pixels are present, keep their hue instead of washing them white.
    saturated = [p for p in bright_core if max(p) - min(p) >= 45]
    if len(saturated) >= max(4, len(bright_core) * 0.25):
        saturated = sorted(saturated, key=lambda p: (max(p) - min(p), sum(p)))
        saturated = saturated[int(len(saturated) * 0.35):]
        return median_color(saturated)

    return median_color(bright_core)


def apply_text_fill(
    run,
    color: tuple[int, int, int],
    gradient: tuple[tuple[int, int, int], tuple[int, int, int]] | None = None,
) -> None:
    if not gradient:
        run.font.color.rgb = RGBColor(*color)
        return

    rpr = run._r.get_or_add_rPr()
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill", "a:grpFill"):
        el = rpr.find(qn(tag))
        if el is not None:
            rpr.remove(el)

    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("rotWithShape", "1")
    gs_lst = OxmlElement("a:gsLst")
    for pos, stop_color in (("0", gradient[0]), ("100000", gradient[1])):
        gs = OxmlElement("a:gs")
        gs.set("pos", pos)
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", color_hex(stop_color))
        gs.append(srgb)
        gs_lst.append(gs)
    lin = OxmlElement("a:lin")
    lin.set("ang", "0")
    lin.set("scaled", "1")
    grad_fill.append(gs_lst)
    grad_fill.append(lin)
    rpr.append(grad_fill)


def parse_rgb(value) -> tuple[int, int, int] | None:
    if value is None:
        return None
    if isinstance(value, str):
        raw = value.strip().lstrip("#")
        if len(raw) == 6:
            try:
                return tuple(int(raw[i:i + 2], 16) for i in (0, 2, 4))
            except ValueError:
                return None
        return None
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            return tuple(max(0, min(255, int(round(float(v))))) for v in value[:3])
        except (TypeError, ValueError):
            return None
    return None


def parse_gradient(value) -> tuple[tuple[int, int, int], tuple[int, int, int]] | None:
    if value is None:
        return None
    if isinstance(value, dict):
        start = parse_rgb(value.get("start") or value.get("left") or value.get("from"))
        end = parse_rgb(value.get("end") or value.get("right") or value.get("to"))
    elif isinstance(value, (list, tuple)) and len(value) >= 2:
        start = parse_rgb(value[0])
        end = parse_rgb(value[1])
    else:
        return None
    if start and end:
        return (start, end)
    return None


def load_style_overrides(path: Path | None) -> dict:
    """Load optional multimodal style refinements keyed by page + region index."""
    if path is None:
        return {}
    data = json.loads(path.read_text(encoding="utf-8"))
    overrides = {}

    if isinstance(data, dict):
        for key, value in data.items():
            if isinstance(key, str) and ":" in key and isinstance(value, dict):
                page, index = key.split(":", 1)
                if page.isdigit() and index.isdigit():
                    overrides[(int(page), int(index))] = value
        for slide in data.get("slides", []):
            page = int(slide.get("page", 0) or 0)
            for item in slide.get("regions", []):
                if "index" in item:
                    overrides[(page, int(item["index"]))] = item
    elif isinstance(data, list):
        for item in data:
            if isinstance(item, dict) and "page" in item and "index" in item:
                overrides[(int(item["page"]), int(item["index"]))] = item
    return overrides


def apply_style_override(
    *,
    override: dict | None,
    base_color: tuple[int, int, int],
    base_gradient: tuple[tuple[int, int, int], tuple[int, int, int]] | None,
    base_size: float,
    base_font: str,
) -> tuple[tuple[int, int, int], tuple[tuple[int, int, int], tuple[int, int, int]] | None, float, str]:
    if not override:
        return base_color, base_gradient, base_size, base_font

    color = parse_rgb(override.get("color") or override.get("rgb")) or base_color
    gradient = parse_gradient(override.get("gradient"))
    fill = override.get("fill")
    if gradient is None and isinstance(fill, dict):
        gradient = parse_gradient(fill.get("gradient") or fill)
        color = parse_rgb(fill.get("color") or fill.get("rgb")) or color
    if gradient is None and override.get("gradient") is False:
        gradient = None
    elif gradient is None:
        gradient = base_gradient

    size = base_size
    if "font_size" in override:
        try:
            size = float(override["font_size"])
        except (TypeError, ValueError):
            size = base_size
    elif "font_size_delta" in override:
        try:
            size = base_size + float(override["font_size_delta"])
        except (TypeError, ValueError):
            size = base_size
    elif "font_size_scale" in override:
        try:
            size = base_size * float(override["font_size_scale"])
        except (TypeError, ValueError):
            size = base_size
    size = max(6.0, min(56.0, round(size * 2) / 2))

    font = str(override.get("font_name") or override.get("font") or base_font).strip() or base_font
    return color, gradient, size, font


def estimate_font_size(height_px: float, text: str, image_h: float) -> float:
    px_per_pt = image_h / (SLIDE_H_EMU / 12700.0)
    raw_pt = height_px / px_per_pt / 1.10
    size = max(7.0, min(42.0, round(raw_pt * 2) / 2))
    if len(text) >= 10:
        size *= 0.94
    return size


def safe_text_box(
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    image_w: float,
    image_h: float,
):
    single_line = "\n" not in text
    long_text = len(text) >= 8 or any(ch in text for ch in "/+&|ABCDEFGHIJKLMNOPQRSTUVWXYZ")
    width_scale = 1.36 if single_line and long_text else 1.26 if single_line else 1.14
    height_scale = 1.46 if single_line else 1.24
    new_w = w * width_scale
    new_h = h * height_scale
    new_x = x - (new_w - w) / 2
    new_y = y - (new_h - h) / 2
    if single_line and long_text:
        size *= 0.96
    if single_line and size >= 28:
        size *= 0.97
    if new_x < 0:
        new_x = 0
    if new_y < 0:
        new_y = 0
    if new_x + new_w > image_w:
        new_x = max(0, image_w - new_w)
    if new_y + new_h > image_h:
        new_y = max(0, image_h - new_h)
    return new_x, new_y, new_w, new_h, size


def sample_text_paint(orig: Image.Image, clean: Image.Image, bbox: list[float]) -> dict:
    x0, y0, x1, y1 = [int(round(v)) for v in bbox]
    pad_x = max(2, int((x1 - x0) * 0.08))
    pad_y = max(2, int((y1 - y0) * 0.18))
    x0 = max(0, min(orig.width - 1, x0 - pad_x))
    y0 = max(0, min(orig.height - 1, y0 - pad_y))
    x1 = max(x0 + 1, min(orig.width, x1 + pad_x))
    y1 = max(y0 + 1, min(orig.height, y1 + pad_y))

    orig_crop = orig.crop((x0, y0, x1, y1))
    clean_crop = clean.crop((x0, y0, x1, y1))
    max_w, max_h = 180, 80
    if orig_crop.width > max_w or orig_crop.height > max_h:
        scale = min(max_w / orig_crop.width, max_h / orig_crop.height)
        new_size = (max(1, int(orig_crop.width * scale)), max(1, int(orig_crop.height * scale)))
        orig_crop = orig_crop.resize(new_size)
        clean_crop = clean_crop.resize(new_size)

    candidates: list[tuple[int, int, int, int]] = []
    clean_pixels = [p[:3] for p in clean_crop.getdata()]
    bg_brightness = statistics.median(sum(p) for p in clean_pixels) if clean_pixels else 0
    crop_w = orig_crop.width
    for idx, (op, cp) in enumerate(zip(orig_crop.getdata(), clean_crop.getdata())):
        r, g, b = op[:3]
        cr, cg, cb = cp[:3]
        diff = abs(r - cr) + abs(g - cg) + abs(b - cb)
        brightness = r + g + b
        saturation = max(r, g, b) - min(r, g, b)
        # Text is what changed between original and no-text background.
        # Filter out tiny antialias noise and very dark background edges.
        px = idx % crop_w
        if diff >= 34 and brightness >= 78:
            candidates.append((px, r, g, b))
        elif diff >= 54 and saturation >= 18:
            candidates.append((px, r, g, b))

    if len(candidates) < 6:
        # Fallback: prefer bright/saturated pixels from the original crop,
        # because this deck uses mostly white/cyan text over dark backgrounds.
        pixels = [(idx % crop_w, *p[:3]) for idx, p in enumerate(orig_crop.getdata())]
        candidates = [
            p for p in pixels
            if sum(p[1:]) >= 210 or (max(p[1:]) - min(p[1:]) >= 35 and sum(p[1:]) >= 120)
        ]
    if not candidates:
        return {"color": (235, 246, 255), "gradient": None}

    rgb_candidates = [p[1:] for p in candidates]
    base_color = brighten_text_color(core_ink_color(rgb_candidates), bg_brightness)

    gradient = None
    if len(candidates) >= 10:
        by_x = sorted(candidates, key=lambda p: p[0])
        edge_n = max(5, int(len(by_x) * 0.30))
        left_color = brighten_text_color(core_ink_color([p[1:] for p in by_x[:edge_n]]), bg_brightness)
        right_color = brighten_text_color(core_ink_color([p[1:] for p in by_x[-edge_n:]]), bg_brightness)
        if color_distance(left_color, right_color) >= 18:
            gradient = (left_color, right_color)
        elif sum(base_color) >= 560 and max(base_color) - min(base_color) < 46 and len(candidates) >= 18:
            # Many small white labels in the source have subtle blue-white
            # lighting. Make it visible without inventing a strong hue shift.
            gradient = (
                mix_color(base_color, (210, 232, 255), 0.34),
                mix_color(base_color, (255, 255, 255), 0.24),
            )

    return {"color": base_color, "gradient": gradient}


def module_exists(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def probe_import(command: str) -> tuple[bool, str]:
    try:
        proc = subprocess.run(
            [sys.executable, "-c", command],
            capture_output=True,
            text=True,
            timeout=45,
            check=False,
        )
    except Exception as exc:
        return (False, str(exc))
    detail = (proc.stdout or proc.stderr or "").strip()
    return (proc.returncode == 0, detail)


def can_run_paddleocr() -> bool:
    if not (module_exists("paddleocr") and module_exists("paddle")):
        return False
    ok, _detail = probe_import("import paddle, paddleocr; print(paddle.__version__)")
    return ok


def can_run_rapidocr() -> bool:
    if not (module_exists("rapidocr") or module_exists("rapidocr_onnxruntime")):
        return False
    ok, _detail = probe_import(
        "from PIL import Image\n"
        "import cv2\n"
        "try:\n"
        "    from rapidocr import RapidOCR\n"
        "except ImportError:\n"
        "    from rapidocr_onnxruntime import RapidOCR"
    )
    return ok


def apply_low_memory_runtime_defaults() -> None:
    for key in ("OMP_NUM_THREADS", "OPENBLAS_NUM_THREADS", "MKL_NUM_THREADS", "NUMEXPR_NUM_THREADS"):
        os.environ.setdefault(key, "1")
    try:
        import cv2

        cv2.setNumThreads(1)
    except Exception:
        pass


def choose_ocr_backend(requested: str) -> str:
    if requested != "auto":
        return requested
    if can_run_paddleocr():
        return "paddleocr"
    if can_run_rapidocr():
        return "rapidocr"
    raise RuntimeError(
        "No OCR backend is available. Install PaddleOCR or RapidOCR, for example: "
        "pip install paddlepaddle paddleocr opencv-python-headless "
        "or pip install rapidocr_onnxruntime opencv-python-headless"
    )


def normalize_regions(image_path: Path, width: float, height: float, rows: list) -> dict:
    regions = []
    for row in rows:
        text = ""
        confidence = 0.0
        box = None
        if isinstance(row, dict):
            text = str(row.get("text") or row.get("rec_text") or row.get("label") or "")
            confidence = float(row.get("confidence") or row.get("score") or row.get("rec_score") or 0.0)
            box = row.get("bbox") or row.get("box") or row.get("dt_box") or row.get("points")
        elif isinstance(row, (list, tuple)) and len(row) >= 3:
            box, text, confidence = row[0], str(row[1]), float(row[2] or 0.0)
        if not text.strip() or box is None:
            continue
        if len(box) == 4 and all(isinstance(v, (int, float)) for v in box):
            x0, y0, x1, y1 = [float(v) for v in box]
        else:
            points = [(float(p[0]), float(p[1])) for p in box]
            xs = [p[0] for p in points]
            ys = [p[1] for p in points]
            x0, y0, x1, y1 = min(xs), min(ys), max(xs), max(ys)
        x0 = max(0.0, min(width, x0))
        x1 = max(0.0, min(width, x1))
        y0 = max(0.0, min(height, y0))
        y1 = max(0.0, min(height, y1))
        if x1 <= x0 or y1 <= y0:
            continue
        regions.append({"text": text.strip(), "confidence": confidence, "bbox": [x0, y0, x1, y1]})
    return {"image": str(image_path), "width": width, "height": height, "regions": regions}


def paddle_result_rows(raw) -> list:
    if raw is None:
        return []
    if hasattr(raw, "to_json"):
        raw = raw.to_json()
    if isinstance(raw, str):
        raw = json.loads(raw)
    if isinstance(raw, dict):
        texts = raw.get("rec_texts") or raw.get("texts") or []
        scores = raw.get("rec_scores") or raw.get("scores") or [0.0] * len(texts)
        boxes = (
            raw.get("rec_polys")
            or raw.get("dt_polys")
            or raw.get("rec_boxes")
            or raw.get("boxes")
            or []
        )
        if texts and boxes:
            return list(zip(boxes, texts, scores))
        for key in ("results", "regions", "res"):
            if key in raw:
                return paddle_result_rows(raw[key])
        return []
    if isinstance(raw, (list, tuple)) and len(raw) == 1:
        return paddle_result_rows(raw[0])
    rows = []
    if isinstance(raw, (list, tuple)):
        for item in raw:
            if isinstance(item, dict):
                rows.extend(paddle_result_rows(item))
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                if len(item) >= 3:
                    rows.append(item)
                elif isinstance(item[1], (list, tuple)) and len(item[1]) >= 2:
                    rows.append((item[0], item[1][0], item[1][1]))
                else:
                    rows.extend(paddle_result_rows(item))
    return rows


def run_paddleocr(image_path: Path, width: float, height: float) -> dict:
    apply_low_memory_runtime_defaults()
    from paddleocr import PaddleOCR

    init_attempts = [
        {"lang": "ch", "use_angle_cls": True, "show_log": False, "use_gpu": False},
        {"lang": "ch", "use_textline_orientation": True},
        {"lang": "ch"},
    ]
    last_error = None
    engine = None
    for kwargs in init_attempts:
        try:
            engine = PaddleOCR(**kwargs)
            break
        except TypeError as exc:
            last_error = exc
    if engine is None:
        raise RuntimeError(f"Failed to initialize PaddleOCR: {last_error}")

    if hasattr(engine, "ocr"):
        try:
            raw = engine.ocr(str(image_path), cls=True)
        except TypeError:
            raw = engine.ocr(str(image_path))
    elif hasattr(engine, "predict"):
        raw = engine.predict(str(image_path))
    else:
        raise RuntimeError("Unsupported PaddleOCR API: missing ocr() or predict()")
    return normalize_regions(image_path, width, height, paddle_result_rows(raw))


def run_rapidocr(image_path: Path, width: float, height: float) -> dict:
    apply_low_memory_runtime_defaults()
    try:
        from rapidocr import RapidOCR
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR

    engine = RapidOCR()
    raw = engine(str(image_path))
    if isinstance(raw, tuple):
        raw = raw[0]

    rows = []
    if raw is None:
        rows = []
    elif hasattr(raw, "boxes") and hasattr(raw, "txts"):
        scores = getattr(raw, "scores", [0.0] * len(raw.txts))
        rows = list(zip(raw.boxes, raw.txts, scores))
    elif hasattr(raw, "to_json"):
        data = raw.to_json()
        if isinstance(data, str):
            data = json.loads(data)
        rows = data.get("results") or data.get("regions") or []
    else:
        rows = list(raw)
    return normalize_regions(image_path, width, height, rows)


def run_ocr(
    image_path: Path,
    out_dir: Path,
    backend: str,
    ocr_json_dir: Path | None,
) -> dict:
    image = Image.open(image_path)
    width, height = image.size
    image.close()

    cache_dir = out_dir / "ocr" / backend
    if backend == "json":
        if ocr_json_dir is None:
            raise RuntimeError("--ocr-json-dir is required when --ocr-backend=json")
        cache_dir = ocr_json_dir

    cache_dir.mkdir(parents=True, exist_ok=True)
    out_json = cache_dir / f"{image_path.stem}.json"
    if out_json.exists():
        print(f"Using cached OCR {backend}/{out_json.name}", flush=True)
        return json.loads(out_json.read_text(encoding="utf-8"))
    if backend == "json":
        raise FileNotFoundError(f"Missing precomputed OCR JSON: {out_json}")

    print(f"OCR {image_path.name} via {backend}", flush=True)
    if backend == "paddleocr":
        result = run_paddleocr(image_path, width, height)
    elif backend == "rapidocr":
        result = run_rapidocr(image_path, width, height)
    else:
        raise RuntimeError(f"Unsupported OCR backend: {backend}")

    out_json.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Cached OCR {out_json}", flush=True)
    return result


def build(
    base: Path,
    output_name: str,
    ocr_backend: str = "auto",
    font_name: str | None = None,
    ocr_json_dir: Path | None = None,
    style_overrides_path: Path | None = None,
) -> None:
    ppt_dir = base / "ppt"
    clean_dir = base / "ppt-clean"
    out_dir = base / "ppt-editable"
    out_pptx = out_dir / output_name
    text_layers = out_dir / f"{Path(output_name).stem}_text_layers.json"
    backend = choose_ocr_backend(ocr_backend)
    font_name = font_name or default_font_name()
    style_overrides = load_style_overrides(style_overrides_path)

    print(f"OCR backend: {backend}", flush=True)
    print(f"Font: {font_name}", flush=True)
    if style_overrides:
        print(f"Style overrides: {len(style_overrides)} regions", flush=True)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]
    all_layers = {"slides": []}

    for page, image_path in enumerate(sorted(ppt_dir.glob("*.png")), 1):
        print(f"Slide {page}: {image_path.name}", flush=True)
        clean_path = clean_dir / f"{image_path.stem}_clean.png"
        if not clean_path.exists():
            raise FileNotFoundError(f"Missing clean background: {clean_path}")
        orig = Image.open(image_path).convert("RGB")
        clean = Image.open(clean_path).convert("RGB")
        image_w, image_h = orig.size
        if clean.size != orig.size:
            raise ValueError(f"Image size mismatch: {image_path.name} {orig.size} vs {clean_path.name} {clean.size}")
        ocr = run_ocr(image_path, out_dir, backend, ocr_json_dir)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(clean_path),
            Emu(0),
            Emu(0),
            width=Emu(SLIDE_W_EMU),
            height=Emu(SLIDE_H_EMU),
        )

        placed = []
        for region_index, region in enumerate(ocr["regions"]):
            text = region["text"].strip()
            if not text:
                continue
            x0, y0, x1, y1 = region["bbox"]
            w = max(1.0, x1 - x0)
            h = max(1.0, y1 - y0)
            size = estimate_font_size(h, text, image_h)
            x, y, bw, bh, size = safe_text_box(text, x0, y0, w, h, size, image_w, image_h)
            paint = sample_text_paint(orig, clean, region["bbox"])
            color = paint["color"]
            gradient = paint["gradient"]
            style_override = style_overrides.get((page, region_index))
            color, gradient, size, region_font = apply_style_override(
                override=style_override,
                base_color=color,
                base_gradient=gradient,
                base_size=size,
                base_font=font_name,
            )

            tb = slide.shapes.add_textbox(
                Emu(px_to_emu_x(x, image_w)),
                Emu(px_to_emu_y(y, image_h)),
                Emu(px_to_emu_x(bw, image_w)),
                Emu(px_to_emu_y(bh, image_h)),
            )
            tf = tb.text_frame
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            set_ea_font(run, region_font)
            run.font.size = Pt(size)
            apply_text_fill(run, color, gradient)
            placed.append({
                **region,
                "index": region_index,
                "font_size": size,
                "font_name": region_font,
                "color": color,
                "gradient": gradient,
                "style_override": bool(style_override),
            })

        print(f"  placed {len(placed)} text boxes", flush=True)
        all_layers["slides"].append(
            {
                "page": page,
                "image": f"ppt/{image_path.name}",
                "clean": f"ppt-clean/{clean_path.name}",
                "ocr_backend": backend,
                "image_size": [image_w, image_h],
                "regions": placed,
            }
        )

    out_dir.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)
    text_layers.write_text(json.dumps(all_layers, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Saved {out_pptx}", flush=True)
    print(f"Saved {text_layers}", flush=True)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build editable PPT from ppt/*.png + ppt-clean/*_clean.png")
    parser.add_argument("--base", default=".", help="Project directory containing ppt/ and ppt-clean/")
    parser.add_argument("--output", default="editable_ocr_gradient.pptx", help="Output pptx filename")
    parser.add_argument(
        "--ocr-backend",
        default="auto",
        choices=("auto", "paddleocr", "rapidocr", "json"),
        help="OCR backend. auto prefers PaddleOCR, then RapidOCR.",
    )
    parser.add_argument("--ocr-json-dir", default=None, help="Directory of precomputed OCR JSON files for --ocr-backend=json")
    parser.add_argument("--font", default=None, help="PPT font name, e.g. PingFang SC or Microsoft YaHei")
    parser.add_argument(
        "--style-overrides",
        default=None,
        help="Optional multimodal style override JSON keyed by page + region index. It may adjust font size, color, gradient, and font only.",
    )
    args = parser.parse_args()
    build(
        Path(args.base).resolve(),
        args.output,
        ocr_backend=args.ocr_backend,
        font_name=args.font,
        ocr_json_dir=Path(args.ocr_json_dir).resolve() if args.ocr_json_dir else None,
        style_overrides_path=Path(args.style_overrides).resolve() if args.style_overrides else None,
    )
